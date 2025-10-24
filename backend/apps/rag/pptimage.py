from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import base64

output_folder = "./extracted_images"

class PptImageUtil:

  def extract_images_from_group(self, group_shape, image_count):
    for shape in group_shape.shapes:
      if shape.shape_type == 6:  # 6 表示组合形状
        image_count = self.extract_images_from_group(shape, image_count)
      elif shape.shape_type == 13:  # 13 表示图片形状
        image = shape.image
        image_bytes = image.blob
        image_ext = image.ext
        image_filename = f"image_{image_count}.{image_ext}"
        image_path = os.path.join(output_folder, image_filename)
        with open(image_path, 'wb') as f:
          f.write(image_bytes)
          image_count += 1
      elif hasattr(shape, 'fill'):
        try:
          image = shape.fill.picture
          image_bytes = image.blob
          image_ext = image.ext
          image_filename = f"image_{image_count}.{image_ext}"
          image_path = os.path.join(self.output_folder, image_filename)
          with open(image_path, 'wb') as f:
            f.write(image_bytes)
          image_count += 1
        except AttributeError:
          continue
    return image_count
        
  def extract_images(self, ppt_path: str):
    base64_images = []
    image_mime_types = {
        'png': 'image/png',
        'jpg': 'image/jpeg',
        'jpeg': 'image/jpeg',
        'gif': 'image/gif',
        'bmp': 'image/bmp',
        'webp': 'image/webp',
        'tiff': 'image/tiff',
    }
    prs = Presentation(ppt_path)
    image_count = 0
    # 处理幻灯片中的图片
    for slide in prs.slides:
      for shape in slide.shapes:
        if shape.shape_type == 13:
          image = shape.image
          image_bytes = image.blob
          image_ext = image.ext
          # 转换为Base64字符串
          base64_str = base64.b64encode(image_bytes).decode('utf-8')
          # 获取对应的MIME类型，默认使用'application/octet-stream'
          mime_type = image_mime_types.get(image_ext, 'application/octet-stream')   
          # 组合成完整的数据URI
          data_uri = f'data:{mime_type};base64,{base64_str}'
          base64_images.append(data_uri)
          image_count += 1
        elif hasattr(shape, 'fill'):
          try:
            image = shape.fill.picture
            image_bytes = image.blob
            image_ext = image.ext
            # 转换为Base64字符串
            base64_str = base64.b64encode(image_bytes).decode('utf-8')
            # 获取对应的MIME类型，默认使用'application/octet-stream'
            mime_type = image_mime_types.get(image_ext, 'application/octet-stream')   
            # 组合成完整的数据URI
            data_uri = f'data:{mime_type};base64,{base64_str}'
            base64_images.append(data_uri)
            image_count += 1
          except AttributeError:
            continue
            
    return base64_images
  
PPTIMAGEUTIL = PptImageUtil()